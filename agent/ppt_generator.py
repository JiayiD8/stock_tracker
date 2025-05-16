# ppt_generator.py
import os
import re
from io import BytesIO

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR

def create_slide_preview(slide_type, content, ticker=None):
    """Create a styled preview image for a slide based on its type and content"""
    fig, ax = plt.subplots(figsize=(10, 5.625))  # 16:9 aspect ratio

    # Set light background with dark text 
    background_color = '#f0f2f6'  
    fig.patch.set_facecolor(background_color)
    ax.set_facecolor(background_color)
    ax.axis('off')

    # Common styling

    title_color = '#1f2937'  # Dark text for titles
    text_color = '#4b5563'   # Medium dark for content
    accent_color = '#3b82f6'  # Blue accent color
    title_size = 20
    # subtitle_size = 16 # Not used in current implementation
    content_size = 12

    if slide_type == "cover":
        # Cover slide
        ax.text(0.5, 0.6, "AI Report of", color=title_color, fontsize=24, ha='center', weight='bold')
        ax.text(0.5, 0.4, ticker, color=accent_color, fontsize=42, ha='center', weight='bold')

    elif slide_type == "key_takeaways":
        # Key takeaways slide
        # ax.text(0.5, 0.85, "Three Key Takeaways ⚠️", color=title_color, fontsize=title_size, ha='center', weight='bold')
        # takeaways = content.split('\n')
        # takeaways = [t for t in takeaways if t.strip() and t.strip().startswith(('1.', '2.', '3.'))]
        # y_positions = [0.7, 0.5, 0.3]
        # for i, takeaway in enumerate(takeaways[:3]):
        #     if i < len(y_positions):
            
        #         short_takeaway = takeaway[:50] + "..." if len(takeaway) > 50 else takeaway
        #         rect = plt.Rectangle((0.1, y_positions[i]-0.05), 0.8, 0.1, fill=True, 
        #                            color='#dbeafe', alpha=0.5, transform=ax.transAxes)
        #         ax.add_patch(rect)
        #         ax.text(0.5, y_positions[i], short_takeaway, color=text_color, fontsize=content_size,
        #               ha='center', wrap=True)
        # NEW ADDED
        ax.text(0.5, 0.85, "Three Key Takeaways ⚠️", color=title_color, fontsize=title_size, ha='center', weight='bold')
        
        # Extract takeway lines with URLs
        takeaways = re.findall(r"(\d+\.\s*#.*?)(?=\n\d+\.\s*#|\Z)", content, re.DOTALL)
        y_positions = [0.7, 0.5, 0.3]
        
        for i, takeaway in enumerate(takeaways[:3]):
            if i < len(y_positions):
                # Check for URL
                url_match = re.search(r'\[(https?://[^\]]+)\]', takeaway)
                # Clean takeaway text
                takeaway_clean = re.sub(r'\[https?://[^\]]+\]', '', takeaway).strip()
                # Shorten for display
                short_takeaway = takeaway_clean[:50] + "..." if len(takeaway_clean) > 50 else takeaway_clean
                
                rect = plt.Rectangle((0.1, y_positions[i]-0.05), 0.8, 0.1, fill=True, 
                                   color='#dbeafe', alpha=0.5, transform=ax.transAxes)
                ax.add_patch(rect)
                ax.text(0.5, y_positions[i], short_takeaway, color=text_color, fontsize=content_size,
                      ha='center', wrap=True)
                
                 # Add a small indicator if source is available
                if url_match:
                    ax.text(0.85, y_positions[i]-0.04, "📄", color=accent_color, fontsize=content_size-2,
                          ha='center', va='center')
    elif slide_type == "macro":
        ax.text(0.5, 0.85, "Investment Environment and Future Prospects 📊",
              color=title_color, fontsize=title_size, ha='center', weight='bold')
        ax.text(0.5, 0.5, "Market and economic analysis for your investment decisions",
              color=text_color, fontsize=content_size, ha='center')
        # Add a subtle background element
        rect = plt.Rectangle((0.1, 0.4), 0.8, 0.2, fill=True, color='#dbeafe', alpha=0.3, transform=ax.transAxes)
        ax.add_patch(rect)
        
    elif slide_type == "catalysts":
        ax.text(0.5, 0.85, "Catalyst ⏳", color=title_color, fontsize=title_size, ha='center', weight='bold')
        catalysts = [l for l in content.split('\n') if l.strip().startswith(('Catalyst 1', 'Catalyst 2', 'Catalyst 3'))]
        y_positions = [0.7, 0.5, 0.3]
        for i, catalyst in enumerate(catalysts[:3]):
            if i < len(y_positions):
                short_catalyst = catalyst[:50] + "..." if len(catalyst) > 50 else catalyst
                # Add a small highlight for each catalyst
                rect = plt.Rectangle((0.1, y_positions[i]-0.05), 0.8, 0.1, fill=True, 
                                   color='#dbeafe', alpha=0.3, transform=ax.transAxes)
                ax.add_patch(rect)
                ax.text(0.5, y_positions[i], short_catalyst, color=text_color, fontsize=content_size,
                      ha='center', wrap=True)
                
    elif slide_type == "price_analysis":
        ax.text(0.5, 0.85, "Stock Price & Volatility Analysis 📈",
              color=title_color, fontsize=title_size, ha='center', weight='bold')
        ax.text(0.5, 0.5, "Analysis of pricing trends and market volatility",
              color=text_color, fontsize=content_size, ha='center')
        # Add a chart-like element
        ax.plot([0.2, 0.4, 0.3, 0.5, 0.7, 0.6, 0.8], 
               [0.3, 0.4, 0.35, 0.5, 0.45, 0.6, 0.7], 
               color=accent_color, linewidth=3)
        
    elif slide_type == "recommendation":
        ax.text(0.5, 0.85, "Investment Recommendation 💰",
              color=title_color, fontsize=title_size, ha='center', weight='bold')
        if "Position" in content:
            position_line = [l for l in content.split('\n') if "Position" in l]
            if position_line:
                position = position_line[0]
                short_position = position[:50] + "..." if len(position) > 50 else position
                rect = plt.Rectangle((0.15, 0.55), 0.7, 0.1, fill=True, 
                                   color='#dbeafe', alpha=0.5, transform=ax.transAxes)
                ax.add_patch(rect)
                ax.text(0.5, 0.6, short_position, color=title_color, fontsize=content_size,
                      ha='center', wrap=True, weight='bold')
        if "Price Target" in content:
            target_line = [l for l in content.split('\n') if "Price Target" in l]
            if target_line:
                target = target_line[0]
                short_target = target[:50] + "..." if len(target) > 50 else target

    buf = BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png', bbox_inches='tight', dpi=100, facecolor='black')
    buf.seek(0)
    plt.close(fig)
    return buf.getvalue()


def create_slide_previews(ticker, financial_report):
    """Create preview images for all slides in the presentation"""
    slides = []
    slides.append(create_slide_preview("cover", None, ticker))
    sections = re.split(r'Section \d+: ', financial_report)[1:]
    if len(sections) > 0:
        slides.append(create_slide_preview("key_takeaways", sections[0].strip()))
    if len(sections) > 1:
        slides.append(create_slide_preview("macro", sections[1].strip()))
    if len(sections) > 2:
        slides.append(create_slide_preview("catalysts", sections[2].strip()))
    if len(sections) > 3:
        slides.append(create_slide_preview("price_analysis", sections[3].strip()))
    if len(sections) > 4:
        slides.append(create_slide_preview("recommendation", sections[4].strip()))
    return slides


def create_ppt(ticker, financial_report, status_text):
    """
    Creates a PowerPoint presentation from a financial report.
    status_text is a Streamlit UI element to update progress.
    """
    status_text.text("Creating PowerPoint presentation...")

    sections = re.split(r'Section \d+: ', financial_report)[1:]

    def clean_footer(text):
        footer = "================================================================================\n🔍 **End of Report** | Generated by AI"
        return text.rsplit(footer, 1)[0].strip() if footer in text else text.strip()

    three_key_takeaways = sections[0].strip() if len(sections) > 0 else ""
    financial_situation_prospects = sections[1].strip() if len(sections) > 1 else ""
    market_catalysts = sections[2].strip() if len(sections) > 2 else ""
    stock_price_volatility = sections[3].strip() if len(sections) > 3 else ""
    investment_recommendation = clean_footer(sections[4]) if len(sections) > 4 else ""

    status_text.text("Building presentation slides...")
    ppt = Presentation()

    # Common font styles
    TEXT_COLOR = RGBColor(31, 41, 55)  # #1f2937 - dark text
    SUBTITLE_COLOR = RGBColor(75, 85, 99)  # #4b5563 - medium dark text
    ACCENT_COLOR = RGBColor(59, 130, 246)  # #3b82f6 - blue accent
    # HIGHLIGHT_BG_COLOR = RGBColor(219, 234, 254)  # #dbeafe - very light blue highlight
    
    SUBTITLE_FONT = "Arial"
    CONTENT_FONT = "Calibri"
    MAIN_TITLE_FONT = "Georgia"
    SLIDE_BACKGROUND_COLOR = RGBColor(240, 242, 246)  # #f0f2f6 - light blue-gray


    def set_slide_background(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = SLIDE_BACKGROUND_COLOR

    def add_main_title(slide, text):
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
        tf = title_box.text_frame
        tf.text = text
        p = tf.paragraphs[0]
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = TEXT_COLOR
        p.font.name = MAIN_TITLE_FONT
        p.alignment = PP_ALIGN.CENTER
        return title_box

    # ---- Cover Slide ----
    slide_layout = ppt.slide_layouts[5]  # Blank slide
    slide = ppt.slides.add_slide(slide_layout)
    set_slide_background(slide)

    title_text_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
    title_tf = title_text_box.text_frame
    title_tf.text = "AI Report of"
    p = title_tf.paragraphs[0]
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = TEXT_COLOR
    p.font.name = MAIN_TITLE_FONT; p.alignment = PP_ALIGN.CENTER

    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3), Inches(4.2), Inches(7), Inches(4.2))
    line.line.color.rgb = ACCENT_COLOR
    line.line.width = Pt(3)
    
    ticker_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(2))
    ticker_tf = ticker_box.text_frame
    ticker_tf.text = ticker
    p = ticker_tf.paragraphs[0]
    p.font.size = Pt(72); p.font.bold = True; p.font.color.rgb = TEXT_COLOR
    p.font.name = SUBTITLE_FONT; p.alignment = PP_ALIGN.CENTER

    # ---- Slide 1: Three Key Takeaways ----
    # NEW ADDED
    slide = ppt.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_main_title(slide, "Three Key Takeaways⚠️")

    # Use a more reliable regex to extract the takeaways
    takeaways_content = re.findall(r'(\d+\.\s*#.*?)(?=\n\d+\.\s*#|\Z)', three_key_takeaways, re.DOTALL)
    box_top_positions = [Inches(1.5), Inches(3.1), Inches(4.7)]  # Adjusted for better spacing

    for i, takeaway in enumerate(takeaways_content[:3]):
        # Extract URL if present
        url_match = re.search(r'\[(https?://[^\]]+)\]', takeaway)
        url = url_match.group(1) if url_match else None
        
        # Clean takeaway text by removing URL
        takeaway_clean = re.sub(r'\[https?://[^\]]+\]', '', takeaway).strip()
        
        # Extract title and content
        match = re.match(r'\d+\.\s*#(.*?):\s*(.*)', takeaway_clean, re.DOTALL)
        if match:
            title = match.group(1).strip()
            content = match.group(2).strip()
            
            # Add subtitle
            subtitle_box = slide.shapes.add_textbox(Inches(1), box_top_positions[i], Inches(8), Inches(0.5))
            tf = subtitle_box.text_frame
            tf.text = title
            p = tf.paragraphs[0]
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = TEXT_COLOR
            p.font.name = SUBTITLE_FONT
            # Add content
            desc_box = slide.shapes.add_textbox(Inches(1), box_top_positions[i] + Inches(0.5), Inches(8), Inches(0.9))
            tf = desc_box.text_frame
            tf.text = content
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_COLOR
            p.font.name = CONTENT_FONT
            
            # Add URL if available
            if url:
                source_box = slide.shapes.add_textbox(Inches(1), box_top_positions[i] + Inches(1.4), Inches(8), Inches(0.3))
                tf = source_box.text_frame
                p = tf.paragraphs[0]
                
                # Add "Source: " text
                r = p.add_run()
                r.text = "Source: "
                r.font.size = Pt(10)
                r.font.italic = True
                r.font.color.rgb = SUBTITLE_COLOR
                
                # Add hyperlinked URL
                r = p.add_run()
                r.text = url
                r.font.size = Pt(10)
                r.font.italic = True
                r.font.color.rgb = ACCENT_COLOR
                
                # Properly add hyperlink with relationship
                try:
                    rId = slide.part.relate_to(url, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink', is_external=True)
                    r._r.get_or_add_rPr().add_hlinkClick(rId)
                except Exception as e:
                    print(f"Error adding hyperlink: {e}")
            
        
    # ---- Slide 2: Investment Environment ----
    slide = ppt.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_main_title(slide, "Investment Environment and Future Prospects📊")

    financial_situation_content = re.findall(r"#(.*?)(?=\n\d+\.\s#|\Z)", financial_situation_prospects, re.DOTALL)
    financial_situation_dict = {}
    pattern = r"^(.*?):\s*(.*)$"
    for item in financial_situation_content:
        match = re.match(pattern, item)
        if match: financial_situation_dict[match.group(1).strip()] = match.group(2).strip()

    left_x, right_x, text_width, box_top = Inches(1), Inches(5.5), Inches(3.5), Inches(1.5) # Adjusted right_x and text_width
    
    if len(financial_situation_dict) >= 1:
        items = list(financial_situation_dict.items())
        # Subtitle 1 (Left)
        subtitle_box = slide.shapes.add_textbox(left_x, box_top, text_width, Inches(0.5))
        tf = subtitle_box.text_frame; tf.text = items[0][0]
        p = tf.paragraphs[0]; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = TEXT_COLOR; p.font.name = SUBTITLE_FONT
        
        desc_box = slide.shapes.add_textbox(left_x, box_top + Inches(0.5), text_width, Inches(4)) # Increased height
        tf = desc_box.text_frame; tf.text = items[0][1]; tf.word_wrap = True
        p = tf.paragraphs[0]; p.font.size = Pt(12); p.font.color.rgb = TEXT_COLOR; p.font.name = CONTENT_FONT

    if len(financial_situation_dict) >= 2:
        items = list(financial_situation_dict.items())
        
        # Subtitle 2 (Right)
        subtitle_box = slide.shapes.add_textbox(right_x, box_top, text_width, Inches(0.5))
        tf = subtitle_box.text_frame; tf.text = items[1][0]
        p = tf.paragraphs[0]; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = TEXT_COLOR; p.font.name = SUBTITLE_FONT

        desc_box = slide.shapes.add_textbox(right_x, box_top + Inches(0.5), text_width, Inches(4)) # Increased height
        tf = desc_box.text_frame; tf.text = items[1][1]; tf.word_wrap = True
        p = tf.paragraphs[0]; p.font.size = Pt(12); p.font.color.rgb = TEXT_COLOR; p.font.name = CONTENT_FONT
        
    # ---- Slide 3: Catalysts ----
    slide = ppt.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_main_title(slide, "Catalyst⏳")

    catalyst_matches = re.findall(r"#Catalyst\s*\d*:\s*(.*?)(?=\n\d+\.\s#|\Z)", market_catalysts, re.DOTALL)
    catalyst_matches = [c.strip() for c in catalyst_matches if c.strip()]
    
    text_left_margin, text_width, box_top_start = Inches(1), Inches(8), Inches(1.5)
    current_top = box_top_start
    line_height = Inches(0.4) # Approximate height per line of text + title
    max_lines_per_slide_heuristic = 10 # Adjust based on desired density
    
    for i, catalyst_text in enumerate(catalyst_matches):
        if i > 0 and i % max_lines_per_slide_heuristic == 0: # Heuristic for new slide
            slide = ppt.slides.add_slide(slide_layout)
            set_slide_background(slide)
            add_main_title(slide, f"Catalyst (Continued)⏳")
            current_top = box_top_start

        text_box = slide.shapes.add_textbox(text_left_margin, current_top, text_width, Inches(0.8)) # Adjust height as needed
        tf = text_box.text_frame; tf.word_wrap = True; tf.clear()
        p = tf.add_paragraph()
        
        run_bold = p.add_run(); run_bold.text = f"Catalyst {i+1}: "
        run_bold.font.bold = True; run_bold.font.size = Pt(14); run_bold.font.color.rgb = TEXT_COLOR; run_bold.font.name = SUBTITLE_FONT
        
        run_regular = p.add_run(); run_regular.text = catalyst_text
        run_regular.font.size = Pt(12); run_regular.font.color.rgb = TEXT_COLOR; run_regular.font.name = CONTENT_FONT
        
        current_top += line_height * (1 + len(catalyst_text) // 80) # Estimate lines for description


    # ---- Slide 4: Stock Price & Volatility Analysis ----
    slide = ppt.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_main_title(slide, "Stock Price & Volatility Analysis📈")

    price_volatility_content = re.findall(r"#(.*?)(?=\n\d+\.\s#|\Z)", stock_price_volatility, re.DOTALL)
    price_volatility_dict = {}
    for item in price_volatility_content:
        match = re.match(pattern, item)
        if match: price_volatility_dict[match.group(1).strip()] = match.group(2).strip()

    left_x, right_x = Inches(0.5), Inches(5.2) # Adjusted for slightly wider content
    upper_text_width, bottom_text_width = Inches(4.5), Inches(9) # Adjusted
    upper_box_top, bottom_box_top = Inches(1.5), Inches(4.0) # Adjusted bottom_box_top

    items = list(price_volatility_dict.items())
    if len(items) >= 1: # Stock Price Analysis (Upper-Left)
        
        subtitle_box = slide.shapes.add_textbox(left_x, upper_box_top, upper_text_width, Inches(0.5))
        tf = subtitle_box.text_frame; tf.text = items[0][0]
        p = tf.paragraphs[0]; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = TEXT_COLOR; p.font.name = SUBTITLE_FONT
        desc_box = slide.shapes.add_textbox(left_x, upper_box_top + Inches(0.4), upper_text_width, Inches(1.8))
        tf = desc_box.text_frame; tf.text = items[0][1]; tf.word_wrap = True
        p = tf.paragraphs[0]; p.font.size = Pt(11); p.font.color.rgb = TEXT_COLOR; p.font.name = CONTENT_FONT

    if len(items) >= 2: # Volatility Analysis (Upper-Right)
        
        subtitle_box = slide.shapes.add_textbox(right_x, upper_box_top, upper_text_width, Inches(0.5))
        tf = subtitle_box.text_frame; tf.text = items[1][0]
        p = tf.paragraphs[0]; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = TEXT_COLOR; p.font.name = SUBTITLE_FONT
        desc_box = slide.shapes.add_textbox(right_x, upper_box_top + Inches(0.4), upper_text_width, Inches(1.8))
        tf = desc_box.text_frame; tf.text = items[1][1]; tf.word_wrap = True
        p = tf.paragraphs[0]; p.font.size = Pt(11); p.font.color.rgb = TEXT_COLOR; p.font.name = CONTENT_FONT
        
    if len(items) >= 3: # What They Reflect (Bottom Full-Width)
        
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), bottom_box_top, bottom_text_width, Inches(0.5))
        tf = subtitle_box.text_frame; tf.text = items[2][0]
        p = tf.paragraphs[0]; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = TEXT_COLOR; p.font.name = SUBTITLE_FONT
        desc_box = slide.shapes.add_textbox(Inches(0.5), bottom_box_top + Inches(0.4), bottom_text_width, Inches(2.5))
        tf = desc_box.text_frame; tf.text = items[2][1]; tf.word_wrap = True
        p = tf.paragraphs[0]; p.font.size = Pt(11); p.font.color.rgb = TEXT_COLOR; p.font.name = CONTENT_FONT

    # ---- Slide 5: Investment Recommendation ----
    slide = ppt.slides.add_slide(slide_layout)
    set_slide_background(slide)
    add_main_title(slide, "Investment Recommendation💰")

    recommendation_content = re.findall(r"#(.*?)(?=\n\d+\.\s#|\Z)", investment_recommendation, re.DOTALL)
    recommendation_dict = {}
    for item in recommendation_content:
        match = re.match(pattern, item)
        if match: recommendation_dict[match.group(1).strip()] = match.group(2).strip()

    text_left_margin, text_width, current_top = Inches(1), Inches(8), Inches(1.5) # Start below title
    
    for i, (subtitle_text, desc_text) in enumerate(recommendation_dict.items()):
        # Estimate lines for spacing
        approx_lines_desc = max(1, len(desc_text) // 70) # Chars per line heuristic
        box_height_desc = Inches(0.3 * approx_lines_desc)
              
        subtitle_box = slide.shapes.add_textbox(text_left_margin, current_top, text_width, Inches(0.4))
        tf = subtitle_box.text_frame; tf.text = subtitle_text
        p = tf.paragraphs[0]; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = TEXT_COLOR; p.font.name = SUBTITLE_FONT
        current_top += Inches(0.4)

        desc_box = slide.shapes.add_textbox(text_left_margin, current_top, text_width, box_height_desc + Inches(0.2))
        tf = desc_box.text_frame; tf.text = desc_text; tf.word_wrap = True
        p = tf.paragraphs[0]; p.font.size = Pt(12); p.font.color.rgb = TEXT_COLOR; p.font.name = CONTENT_FONT
        current_top += box_height_desc + Inches(0.3) # Add spacing for next section

    ppt_file_path = f"{ticker}_financial_report.pptx"
    ppt.save(ppt_file_path)
    status_text.text(f"PowerPoint presentation created: {ppt_file_path}")
    return ppt_file_path

def convert_ppt_to_images(ppt_file):
    try:
        slides = []
        prs = Presentation(ppt_file)
        for i, _ in enumerate(prs.slides):
            plt.figure(figsize=(10, 5.625)) # 16:9 aspect ratio
            background_color = '#f0f2f6'  # Light blue-gray background
            text_color = '#1f2937'  # Dark text
            
            plt.text(0.5, 0.5, f"Slide {i+1}\n(Preview)", fontsize=30, ha='center', va='center', color='white', backgroundcolor='black')
            plt.axis('off')
            plt.tight_layout(pad=0)
            
            buf = BytesIO()
            plt.savefig(buf, format='png', facecolor='black', bbox_inches='tight', pad_inches=0)
            buf.seek(0)
            plt.close()
            slides.append(buf.getvalue())
        return slides
    except Exception as e:
        print(f"Error converting PPT to images (placeholder): {str(e)}")
        # Fallback: create very simple placeholders if Presentation object fails
        slide_count = 0
        try:
            prs_check = Presentation(ppt_file)
            slide_count = len(prs_check.slides)
        except: # If even opening fails
            slide_count = 5 # Assume 5 slides for basic placeholders

        slides = []
        for i in range(slide_count):
            fig, ax = plt.subplots(figsize=(10, 5.625))
            background_color = '#f0f2f6'  # Light blue-gray background
            text_color = '#1f2937'  # Dark text
            
            fig.patch.set_facecolor(background_color)
            ax.set_facecolor(background_color)
            ax.text(0.5, 0.5, f"Slide {i + 1}\n(Basic Preview)", color=text_color, fontsize=20, ha='center', va='center')
            ax.axis('off')
            buf = BytesIO()
            plt.savefig(buf, format='png', facecolor=background_color)
            buf.seek(0)
            plt.close(fig)
            slides.append(buf.getvalue())
        return slides
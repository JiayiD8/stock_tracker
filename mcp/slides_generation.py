from typing import Any
import os
import re
from mcp.server.fastmcp import FastMCP

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR

mcp = FastMCP("slides_generator")

DEFAULT_OUTPUT_DIR = "/Users/dongjiayi/Desktop/Presentations"

TEXT_COLOR = RGBColor(31, 41, 55)      
SUBTITLE_COLOR = RGBColor(75, 85, 99)  
ACCENT_COLOR = RGBColor(59, 130, 246)  
SUBTITLE_FONT = "Arial"
CONTENT_FONT = "Calibri"
MAIN_TITLE_FONT = "Georgia"
SLIDE_BACKGROUND_COLOR = RGBColor(240, 242, 246)  

def set_slide_background(slide):
    """Set the background color for a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = SLIDE_BACKGROUND_COLOR

def add_main_title(slide, text):
    """Add a main title to a slide."""
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(1))
    tf = title_box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = TEXT_COLOR
    p.font.name = MAIN_TITLE_FONT
    p.alignment = PP_ALIGN.CENTER
    return title_box

def create_title_slide(ppt, title):
    """Create a title slide with the given title.
    """
    slide_layout = ppt.slide_layouts[5]  
    slide = ppt.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Add title
    title_text_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
    title_tf = title_text_box.text_frame
    title_tf.text = title
    p = title_tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_COLOR
    p.font.name = MAIN_TITLE_FONT
    p.alignment = PP_ALIGN.CENTER

    # Add decorative line
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3), Inches(4.2), Inches(7), Inches(4.2))
    line.line.color.rgb = ACCENT_COLOR
    line.line.width = Pt(3)
    
    return slide

@mcp.tool()
async def create_presentation(content: str, output_dir: str = DEFAULT_OUTPUT_DIR, title: str = "Presentation") -> str:
    """Create a PowerPoint presentation with the provided content.
    """
    os.makedirs(output_dir, exist_ok=True)
    
    ppt = Presentation()
    create_title_slide(ppt, title)
    
    slide_contents = [s.strip() for s in content.split('---') if s.strip()]
    
    for slide_content in slide_contents:
        lines = slide_content.split('\n')
        slide_title = lines[0].strip()
        slide_body = '\n'.join(lines[1:]) if len(lines) > 1 else ""
        
        slide_layout = ppt.slide_layouts[5]  
        slide = ppt.slides.add_slide(slide_layout)
        set_slide_background(slide)
        
        # Add title
        add_main_title(slide, slide_title)
        
        # Add content if there is any
        if slide_body:
            content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
            tf = content_box.text_frame
            tf.text = slide_body
            tf.word_wrap = True
            
            for p in tf.paragraphs:
                p.font.size = Pt(16)
                p.font.color.rgb = TEXT_COLOR
                p.font.name = CONTENT_FONT
    
    # Save the presentation
    output_filename = f"{title.replace(' ', '_')}.pptx"
    output_path = os.path.join(output_dir, output_filename)
    ppt.save(output_path)
    
    return f"Presentation saved to: {output_path}"

@mcp.tool()
async def create_structured_presentation(content: str, output_dir: str = DEFAULT_OUTPUT_DIR, title: str = "Presentation") -> str:
    """Create a structured PowerPoint presentation with sections and subsections.
    """
    os.makedirs(output_dir, exist_ok=True)
    sections = re.split(r'Section \d+: ', content)

    if not sections[0].strip():
        sections = sections[1:]
    
    ppt = Presentation()
    
    create_title_slide(ppt, title)
    
    for i, section_content in enumerate(sections):
        if not section_content.strip():
            continue
            
        subsections = section_content.split("\n#")
        first_part = subsections[0].strip()
        title_lines = first_part.split('\n', 1)
        section_title = title_lines[0].strip()
        intro_content = title_lines[1] if len(title_lines) > 1 else ""
        
        # Create main section slide
        slide_layout = ppt.slide_layouts[5]  
        slide = ppt.slides.add_slide(slide_layout)
        set_slide_background(slide)
        add_main_title(slide, section_title)
        
        # Add intro content if any
        if intro_content:
            content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
            tf = content_box.text_frame
            tf.text = intro_content.strip()
            tf.word_wrap = True
            
            for p in tf.paragraphs:
                p.font.size = Pt(16)
                p.font.color.rgb = TEXT_COLOR
                p.font.name = CONTENT_FONT
        
        for subsection in subsections[1:]:
            # Extract subsection title and content
            subsection_parts = subsection.split('\n', 1)
            subsection_title = subsection_parts[0].strip()
            subsection_content = subsection_parts[1] if len(subsection_parts) > 1 else ""
            
            # Create a slide for the subsection
            slide = ppt.slides.add_slide(slide_layout)
            set_slide_background(slide)
            add_main_title(slide, subsection_title)
            
            # Add subsection content
            content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
            tf = content_box.text_frame
            tf.text = subsection_content.strip()
            tf.word_wrap = True
            
            for p in tf.paragraphs:
                p.font.size = Pt(16)
                p.font.color.rgb = TEXT_COLOR
                p.font.name = CONTENT_FONT
    
    # Save the presentation
    output_filename = f"{title.replace(' ', '_')}.pptx"
    output_path = os.path.join(output_dir, output_filename)
    ppt.save(output_path)
    
    return f"Structured presentation saved to: {output_path}"

@mcp.tool()
async def create_themed_presentation(content: str, output_dir: str = DEFAULT_OUTPUT_DIR, title: str = "Presentation", theme: str = "default") -> str:
    """Create a PowerPoint presentation with a specific theme.
    """
    # Create the output directory if it doesn't exist
    os.makedirs(output_dir, exist_ok=True)
    
    # Create a new presentation
    ppt = Presentation()
    
    # Define theme colors
    themes = {
        "default": {
            "background": RGBColor(240, 242, 246), 
            "text": RGBColor(31, 41, 55),          
            "accent": RGBColor(59, 130, 246)      
        },
        "dark": {
            "background": RGBColor(31, 41, 55),    
            "text": RGBColor(240, 242, 246),       
            "accent": RGBColor(236, 72, 153)     
        },
        "light": {
            "background": RGBColor(255, 255, 255), 
            "text": RGBColor(31, 41, 55),          
            "accent": RGBColor(16, 185, 129)      
        },
        "corporate": {
            "background": RGBColor(243, 244, 246), 
            "text": RGBColor(17, 24, 39),          
            "accent": RGBColor(37, 99, 235)        
        }
    }
    
    selected_theme = themes.get(theme.lower(), themes["default"])
    
    global SLIDE_BACKGROUND_COLOR, TEXT_COLOR, ACCENT_COLOR
    orig_background, orig_text, orig_accent = SLIDE_BACKGROUND_COLOR, TEXT_COLOR, ACCENT_COLOR
    
    try:
        SLIDE_BACKGROUND_COLOR = selected_theme["background"]
        TEXT_COLOR = selected_theme["text"]
        ACCENT_COLOR = selected_theme["accent"]
        
        create_title_slide(ppt, title)
        slide_contents = [s.strip() for s in content.split('---') if s.strip()]
        
        for slide_content in slide_contents:
            lines = slide_content.split('\n')
            slide_title = lines[0].strip()
            slide_body = '\n'.join(lines[1:]) if len(lines) > 1 else ""
            
            # Create a new slide
            slide_layout = ppt.slide_layouts[5]  
            slide = ppt.slides.add_slide(slide_layout)
            set_slide_background(slide)
            
            # Add title
            add_main_title(slide, slide_title)
            
            # Add content if there is any
            if slide_body:
                content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
                tf = content_box.text_frame
                tf.text = slide_body
                tf.word_wrap = True
                
                for p in tf.paragraphs:
                    p.font.size = Pt(16)
                    p.font.color.rgb = TEXT_COLOR
                    p.font.name = CONTENT_FONT
        
        # Save the presentation
        output_filename = f"{title.replace(' ', '_')}.pptx"
        output_path = os.path.join(output_dir, output_filename)
        ppt.save(output_path)
        
        return f"Themed presentation saved to: {output_path}"
    finally:
        SLIDE_BACKGROUND_COLOR, TEXT_COLOR, ACCENT_COLOR = orig_background, orig_text, orig_accent

if __name__ == "__main__":
    mcp.run(transport='stdio')